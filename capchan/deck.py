"""Output stage — render the deck and the diagnostics.

Two artifacts come out of here:
  - the deck itself: full-resolution slide frames assembled into a PDF (always)
    and optionally a PPTX, one captured frame per slide;
  - the diagnostic plot: the difference series with the gate threshold and the
    captured/too-short segments shaded, so the later tuning pass is a glance, not
    a guess.
"""

from __future__ import annotations

import subprocess
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402


def extract_frame(video, t: float, out_path: Path) -> None:
    """Grab one full-resolution frame at time `t`.

    `t` is a static-segment midpoint, so a fast (keyframe) seek lands well inside
    the held slide even if it is not frame-exact.
    """
    cmd = [
        "ffmpeg", "-ss", f"{t:.3f}", "-i", str(video),
        "-frames:v", "1", "-q:v", "2", "-y", "-v", "error", str(out_path),
    ]
    subprocess.run(cmd, check=True)


def build_pdf(image_paths, out_path: Path) -> None:
    images = [Image.open(p).convert("RGB") for p in image_paths]
    if not images:
        return
    images[0].save(out_path, "PDF", save_all=True, append_images=images[1:])


def build_pptx(image_paths, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    for p in image_paths:
        slide = prs.slides.add_slide(blank)
        with Image.open(p) as im:
            iw, ih = im.size
        scale = min(sw / iw, sh / ih)
        w, h = int(iw * scale), int(ih * scale)
        slide.shapes.add_picture(str(p), int((sw - w) / 2), int((sh - h) / 2), width=w, height=h)
    prs.save(str(out_path))


def plot_diff(series, threshold: float, segments, out_path: Path) -> None:
    fig, ax = plt.subplots(figsize=(14, 4))
    ax.plot(series.times, series.diffs, lw=0.8, color="#333333",
            label="frame-to-frame difference")
    ax.axhline(threshold, color="#c0392b", lw=1.0, ls="--",
               label=f"static threshold = {threshold:.4f}")

    labelled = {"captured": False, "short": False}
    for s in segments:
        if s.passed:
            ax.axvspan(s.start_t, s.end_t, color="#2ecc71", alpha=0.25,
                       label=None if labelled["captured"] else "captured static segment")
            labelled["captured"] = True
        elif s.end_t > s.start_t:  # zero-width (single-frame) blocks would not show
            ax.axvspan(s.start_t, s.end_t, color="#bdc3c7", alpha=0.35,
                       label=None if labelled["short"] else "static but too short")
            labelled["short"] = True

    ax.set_xlabel("time (s)")
    ax.set_ylabel("mean abs diff (0-1)")
    ax.set_title("Difference series and stationarity gate")
    ax.margins(x=0.005)
    ax.legend(loc="upper right", fontsize=8)
    fig.tight_layout()
    fig.savefig(out_path, dpi=110)
    plt.close(fig)
