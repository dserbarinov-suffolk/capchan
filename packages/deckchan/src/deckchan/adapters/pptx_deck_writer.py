from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from deckchan.domain.slide import Slide


class PptxDeckWriter:
    def write(self, slides: list[Slide], path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        sw, sh = int(prs.slide_width), int(prs.slide_height)
        for slide_model in slides:
            slide = prs.slides.add_slide(blank)
            with Image.open(slide_model.representative_frame) as image:
                iw, ih = image.size
            scale = min(sw / iw, sh / ih)
            w, h = int(iw * scale), int(ih * scale)
            slide.shapes.add_picture(
                str(slide_model.representative_frame),
                int((sw - w) / 2),
                int((sh - h) / 2),
                width=w,
                height=h,
            )
        prs.save(str(path))
